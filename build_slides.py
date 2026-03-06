"""
build_slides.py
Generates slides.pptx for the ClarityPay take-home exercise.
Run from the ClarityPay/ project root:
    python3 build_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.enum.dml import MSO_THEME_COLOR
import copy
import os

# ── Dimensions ────────────────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)

# ── Color palette (matches notebook charts) ───────────────────
BLUE   = RGBColor(0x0D, 0x47, 0xA1)   # header bar
BLUE_L = RGBColor(0x1A, 0x6E, 0xD6)   # lighter blue accent
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
RED    = RGBColor(0xEF, 0x53, 0x50)
LGRAY  = RGBColor(0xF5, 0xF5, 0xF5)
MGRAY  = RGBColor(0xE0, 0xE0, 0xE0)
DGRAY  = RGBColor(0x42, 0x42, 0x42)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)

FIGS = 'figs/'


# ─────────────────────────────────────────────────────────────
# HELPER UTILITIES
# ─────────────────────────────────────────────────────────────

def blank_slide(prs):
    """Add a new blank slide (layout 6 = completely blank)."""
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


def fill_slide_bg(slide, color):
    """Fill the entire slide background with a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line=False):
    """Add a filled rectangle shape."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if not line:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                size=18, bold=False, color=DGRAY,
                align=PP_ALIGN.LEFT, wrap=True, italic=False):
    """Add a single-paragraph text box."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_bullet_list(slide, items, left, top, width, height,
                    size=17, color=DGRAY, bullet_char='▸ '):
    """Text box with multiple bullet lines."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = bullet_char + item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        p.space_before = Pt(5)
    return tb


def add_header_bar(slide, title_text, subtitle=None):
    """Dark blue header bar with white title (and optional subtitle)."""
    bar_h = Inches(1.1)
    add_rect(slide, 0, 0, W, bar_h, BLUE)
    tb = slide.shapes.add_textbox(Inches(0.35), Inches(0.1), Inches(12.6), Inches(0.85))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.bold = True
    run.font.size = Pt(27)
    run.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(15)
        r2.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)


def add_image(slide, path, left, top, width, height=None):
    """Insert a PNG image. height=None = auto from aspect ratio."""
    if height:
        slide.shapes.add_picture(path, left, top, width, height)
    else:
        slide.shapes.add_picture(path, left, top, width)


def set_cell(cell, text, size=14, bold=False,
             bg=None, fg=WHITE, align=PP_ALIGN.LEFT):
    """Set text and optional background color on a table cell."""
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = fg
    if bg:
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        solidFill = _make_solid_fill(bg)
        tcPr.append(solidFill)


def _make_solid_fill(rgb):
    """Return an <a:solidFill> XML element for table cell backgrounds."""
    from lxml import etree
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    solidFill = etree.SubElement(etree.Element('dummy'), qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', f'{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}')
    return solidFill


def add_divider(slide, top, color=MGRAY):
    """Thin horizontal rule across most of the slide width."""
    add_rect(slide, Inches(0.3), top, Inches(12.7), Inches(0.02), color)


# ─────────────────────────────────────────────────────────────
# SLIDE 1 — TITLE
# ─────────────────────────────────────────────────────────────
def slide_title(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, BLUE)

    # Decorative accent bar at bottom
    add_rect(slide, 0, H - Inches(0.55), W, Inches(0.55),
             RGBColor(0x1A, 0x5E, 0xC8))

    # Red accent stripe
    add_rect(slide, 0, H - Inches(0.6), Inches(3.5), Inches(0.08), RED)

    # Main title
    tb = slide.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(11.5), Inches(2.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = 'Are We Lending to\nthe Right Borrowers?'
    run.font.bold = True
    run.font.size = Pt(44)
    run.font.color.rgb = WHITE

    # Subtitle
    add_textbox(slide,
                'Consumer Credit Policy Analysis  ·  Lending Club 2007–2020 Q3',
                Inches(0.9), Inches(4.15), Inches(11.5), Inches(0.7),
                size=22, color=RGBColor(0xBB, 0xDE, 0xFB))

    # Company / exercise label
    add_textbox(slide,
                'ClarityPay Take-Home Exercise',
                Inches(0.9), Inches(4.95), Inches(8), Inches(0.5),
                size=17, color=RGBColor(0x90, 0xCA, 0xF9), italic=True)

    # Dataset badge
    add_rect(slide, Inches(0.9), Inches(5.9), Inches(4.2), Inches(0.55),
             RGBColor(0x1A, 0x5E, 0xC8))
    add_textbox(slide, '  1,860,331 completed loans  |  2007–2020 Q3',
                Inches(0.95), Inches(5.9), Inches(4.1), Inches(0.55),
                size=14, color=WHITE)


# ─────────────────────────────────────────────────────────────
# SLIDE 2 — THE PROBLEM
# ─────────────────────────────────────────────────────────────
def slide_problem(prs):
    slide = blank_slide(prs)
    add_header_bar(slide, '1 in 5 Lending Club Loans Were Never Repaid')

    # Left: outcome distribution chart
    add_image(slide, FIGS + 'outcome_distribution.png',
              Inches(0.2), Inches(1.2), Inches(6.2), Inches(5.3))

    # Right: big stat
    add_textbox(slide, '19.5%',
                Inches(7.0), Inches(1.3), Inches(5.5), Inches(2.0),
                size=72, bold=True, color=RED)

    add_textbox(slide, 'charge-off (default) rate',
                Inches(7.0), Inches(3.1), Inches(5.5), Inches(0.5),
                size=19, color=DGRAY)

    add_divider(slide, Inches(3.75))

    add_bullet_list(slide, [
        '362,548 loans written off as uncollectible',
        'Each charge-off = 100% loss of outstanding principal',
        'Industry average for consumer installment: 2–5%',
        'LC\'s platform attracted near-prime / subprime borrowers',
        'Goal: identify high-risk borrowers before we lend',
    ],
        left=Inches(6.8), top=Inches(3.9),
        width=Inches(6.0), height=Inches(3.2),
        size=16, color=DGRAY)


# ─────────────────────────────────────────────────────────────
# SLIDE 3 — DATA OVERVIEW
# ─────────────────────────────────────────────────────────────
def slide_data(prs):
    slide = blank_slide(prs)
    add_header_bar(slide,
                   '1.86 Million Completed Loans Analyzed',
                   subtitle='Scope: Fully Paid or Charged Off only  ·  Origination-time features only  ·  No post-loan data leakage')

    # Left: volume/year chart
    add_image(slide, FIGS + 'volume_bad_rate_by_year.png',
              Inches(0.2), Inches(1.35), Inches(7.5), Inches(5.3))

    # Right: stats block
    stats = [
        ('Total loans analyzed', '1,860,331'),
        ('Fully Paid',           '1,497,783  (80.5%)'),
        ('Charged Off',          '362,548  (19.5%)'),
        ('Date range',           'Jan 2007 – Sep 2020'),
    ]

    top = Inches(1.55)
    for label, value in stats:
        add_textbox(slide, label,
                    Inches(8.1), top, Inches(5.0), Inches(0.38),
                    size=14, color=RGBColor(0x75, 0x75, 0x75))
        add_textbox(slide, value,
                    Inches(8.1), top + Inches(0.34), Inches(5.0), Inches(0.5),
                    size=20, bold=True, color=DGRAY)
        add_divider(slide, top + Inches(0.86), LGRAY)
        top += Inches(1.05)

    # Methodology note
    add_textbox(slide,
                'Methodology: Only loans with a known outcome are included. '
                'In-progress loans (Current, Late, Grace Period) are excluded. '
                'All features used are observable at origination — no post-loan data.',
                Inches(8.1), Inches(6.0), Inches(5.0), Inches(1.2),
                size=11, color=RGBColor(0x90, 0x90, 0x90), italic=True)


# ─────────────────────────────────────────────────────────────
# SLIDE 4 — KEY RISK DRIVERS
# ─────────────────────────────────────────────────────────────
def slide_risk_drivers(prs):
    slide = blank_slide(prs)
    add_header_bar(slide, 'Four Attributes Predict Default Before the Loan Is Made')

    # 2×2 grid of charts
    charts = [
        (FIGS + 'risk_by_fico.png',      'FICO Score at Origination'),
        (FIGS + 'risk_by_grade.png',     'Loan Grade (A–G)'),
        (FIGS + 'risk_by_revol_util.png','Revolving Credit Utilization'),
        (FIGS + 'risk_by_purpose.png',   'Loan Purpose'),
    ]

    positions = [
        (Inches(0.1),  Inches(1.15)),  # top-left
        (Inches(6.72), Inches(1.15)),  # top-right
        (Inches(0.1),  Inches(4.05)),  # bottom-left
        (Inches(6.72), Inches(4.05)),  # bottom-right
    ]

    cw, ch = Inches(6.55), Inches(2.85)

    for (path, label), (left, top) in zip(charts, positions):
        add_image(slide, path, left, top, cw, ch)


# ─────────────────────────────────────────────────────────────
# SLIDE 5 — POLICY RULES TABLE
# ─────────────────────────────────────────────────────────────
def slide_rules(prs):
    slide = blank_slide(prs)
    add_header_bar(slide, 'Six Evidence-Based Decline Rules')

    rows_data = [
        # (Rule, Threshold, Bad Rate if Declined, Rationale)
        ('Rule',             'Threshold',          'Bad Rate if Declined', 'Business Rationale'),
        ('FICO Score',       '< 680',              '24.4%',  'Near-subprime; 2× the average default rate'),
        ('Loan Grade',       'E, F, or G',         '39.9%',  'Lending Club\'s own model flags >28% expected default'),
        ('Revolving Util',   '> 80%',              '21.6%',  'Maxed-out credit = severe financial stress'),
        ('Bankruptcy',       'Any record',         '22.5%',  'Prior bankruptcy = proven inability to repay'),
        ('Credit Inquiries', '≥ 4 in last 6 mo',  '~27%',   'Rapid credit-seeking = financial desperation signal'),
        ('DTI',              '> 35%',              '~29%',   'Monthly debt burden exceeds safe capacity threshold'),
    ]

    col_widths = [Inches(2.1), Inches(2.3), Inches(2.3), Inches(6.0)]
    table = slide.shapes.add_table(
        7, 4,
        Inches(0.3), Inches(1.25),
        sum(col_widths), Inches(5.8)
    ).table

    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for r, row in enumerate(rows_data):
        for c, text in enumerate(row):
            cell = table.cell(r, c)
            if r == 0:
                set_cell(cell, text, size=14, bold=True, bg=BLUE, fg=WHITE,
                         align=PP_ALIGN.CENTER)
            else:
                bg = LGRAY if r % 2 == 0 else WHITE
                txt_color = RED if c == 2 else DGRAY
                bold_col = c in (0, 2)
                set_cell(cell, text, size=13, bold=bold_col,
                         bg=bg, fg=txt_color, align=PP_ALIGN.LEFT)

    # Footnote
    add_textbox(slide,
                'Baseline bad rate: 19.5%  ·  All rules use only origination-time data',
                Inches(0.3), Inches(7.1), Inches(10), Inches(0.32),
                size=11, color=RGBColor(0x90, 0x90, 0x90), italic=True)


# ─────────────────────────────────────────────────────────────
# SLIDE 6 — TRADE-OFF ANALYSIS
# ─────────────────────────────────────────────────────────────
def slide_tradeoff(prs):
    slide = blank_slide(prs)
    add_header_bar(slide,
                   'Recommended Ruleset: FICO < 680  OR  Grade E/F/G  OR  Revol Util > 80%')

    # Left: trade-off scatter
    add_image(slide, FIGS + 'tradeoff_scatter.png',
              Inches(0.1), Inches(1.2), Inches(7.8), Inches(5.8))

    # Right: metrics panel header
    add_rect(slide, Inches(8.2), Inches(1.35), Inches(4.8), Inches(0.55), BLUE)
    add_textbox(slide, '  Recommended Ruleset (R4)',
                Inches(8.2), Inches(1.35), Inches(4.8), Inches(0.55),
                size=15, bold=True, color=WHITE)

    metrics = [
        ('Volume Declined',   '43.3%',  DGRAY),
        ('New Bad Rate',       '15.5%',  GREEN),
        ('Bad Rate Reduction', '−20.5%', GREEN),
        ('Good Loans Lost',    '40.5%',  RED),
    ]

    top = Inches(2.05)
    for label, value, val_color in metrics:
        add_textbox(slide, label,
                    Inches(8.3), top, Inches(3.2), Inches(0.35),
                    size=13, color=RGBColor(0x75, 0x75, 0x75))
        add_textbox(slide, value,
                    Inches(11.0), top, Inches(1.9), Inches(0.35),
                    size=16, bold=True, color=val_color, align=PP_ALIGN.RIGHT)
        add_divider(slide, top + Inches(0.38), LGRAY)
        top += Inches(0.62)

    # Callout
    add_rect(slide, Inches(8.2), Inches(4.6), Inches(4.8), Inches(0.8),
             RGBColor(0xE3, 0xF2, 0xFD))
    add_textbox(slide,
                '"Elbow" of the trade-off curve — meaningful risk\n'
                'reduction without catastrophic volume loss.',
                Inches(8.3), Inches(4.65), Inches(4.6), Inches(0.75),
                size=13, color=BLUE, italic=True)

    # Bottom: ruleset comparison chart
    add_image(slide, FIGS + 'ruleset_comparison.png',
              Inches(8.2), Inches(5.55), Inches(4.8), Inches(1.75))


# ─────────────────────────────────────────────────────────────
# SLIDE 7 — RECOMMENDATIONS
# ─────────────────────────────────────────────────────────────
def slide_recommendations(prs):
    slide = blank_slide(prs)
    add_header_bar(slide, 'Recommendations & Next Steps')

    # Left column header
    add_rect(slide, Inches(0.3), Inches(1.25), Inches(6.0), Inches(0.52), GREEN)
    add_textbox(slide, '  ✓  Implement Now',
                Inches(0.3), Inches(1.25), Inches(6.0), Inches(0.52),
                size=17, bold=True, color=WHITE)

    add_bullet_list(slide, [
        'Decline if FICO < 680\n  → strongest single predictor; 2× baseline bad rate',
        'Decline if Loan Grade E, F, or G\n  → Lending Club\'s own risk model endorses these thresholds',
        'Decline if Revolving Utilization > 80%\n  → captures financial stress independent of credit score',
        'Add Bankruptcy as a hard cutoff (any record)\n  → categorical signal; no exceptions policy',
    ],
        left=Inches(0.3), top=Inches(1.9),
        width=Inches(6.0), height=Inches(4.9),
        size=14, bullet_char='')

    # Right column header
    add_rect(slide, Inches(7.0), Inches(1.25), Inches(6.0), Inches(0.52), BLUE)
    add_textbox(slide, '  →  Future Improvements',
                Inches(7.0), Inches(1.25), Inches(6.0), Inches(0.52),
                size=17, bold=True, color=WHITE)

    add_bullet_list(slide, [
        'Build a logistic regression scorecard\n  → replace hard cutoffs with a continuous risk score; tune threshold at any desired approval rate',
        'Vintage analysis\n  → calibrate rules on recent cohorts (2017–2019) to avoid over-weighting the 2008 crisis',
        'Purpose-specific policies\n  → small business loans (26% bad rate) may need a separate, stricter policy',
        'A/B test the rollout\n  → measure real approval rate impact before full deployment',
    ],
        left=Inches(7.0), top=Inches(1.9),
        width=Inches(6.0), height=Inches(4.9),
        size=14, bullet_char='')

    # Bottom footnote bar
    add_rect(slide, 0, H - Inches(0.42), W, Inches(0.42), LGRAY)
    add_textbox(slide,
                'All rules use only origination-time data  ·  '
                'No post-loan information used  ·  '
                'Analysis based on 1,860,331 Lending Club loans 2007–2020 Q3',
                Inches(0.3), H - Inches(0.4), Inches(12.7), Inches(0.4),
                size=10, color=RGBColor(0x75, 0x75, 0x75), italic=True)


# ─────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    print("Building slides...")
    slide_title(prs);            print("  Slide 1 — Title ✓")
    slide_problem(prs);          print("  Slide 2 — The Problem ✓")
    slide_data(prs);             print("  Slide 3 — Data Overview ✓")
    slide_risk_drivers(prs);     print("  Slide 4 — Key Risk Drivers ✓")
    slide_rules(prs);            print("  Slide 5 — Policy Rules ✓")
    slide_tradeoff(prs);         print("  Slide 6 — Trade-off Analysis ✓")
    slide_recommendations(prs);  print("  Slide 7 — Recommendations ✓")

    out = 'slides.pptx'
    prs.save(out)
    size_kb = os.path.getsize(out) / 1024
    print(f"\nSaved → {out}  ({size_kb:.0f} KB, {len(prs.slides)} slides)")


if __name__ == '__main__':
    main()
